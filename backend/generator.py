import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from .utils import read_file_content
from .ai_client import OllamaClient

class PPTXGenerator:
    def __init__(self, model="gemma3:270m"):
        self.client = OllamaClient(model=model)
        self.prs = Presentation()

    def _add_text(self, text_frame, text, is_title=False):
        """
        Adds text with standard styling.
        """
        p = text_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(32) if is_title else Pt(18)
        p.font.bold = is_title
        p.font.color.rgb = RGBColor(0, 0, 0) # Black text

    def generate_presentation(self, input_files, output_file, image_path=None):
        # 1. Read Content
        full_content = ""
        for file_path in input_files:
            try:
                content = read_file_content(file_path)
                full_content += f"\n--- Content from {os.path.basename(file_path)} ---\n{content}\n"
            except Exception as e:
                print(f"Skipping file {file_path}: {e}")

        if not full_content:
            return "No content found in selected files."

        # 2. Get Structure from AI
        prompt = (
            "Task: Generate a presentation based on the INPUT below.\n"
            "If the input is a command (e.g., 'Make a presentation about X'), ignore the command and just generate slides about X.\n"
            "If the input is text data, summarize it into slides.\n"
            "Use the SAME LANGUAGE as the INPUT.\n"
            "Output format MUST be:\n"
            "Slide: [Title]\n"
            "- [Point]\n"
            "- [Point]\n"
            "\n"
            "INPUT:\n"
            f"{full_content}\n"
            "\n"
            "OUTPUT (Slides only):"
        )
        
        print("Sending request to Ollama...")
        print(f"\n[DEBUG] PROMPT SENT:\n{prompt}\n")
        response_text = self.client.generate_text(prompt)
        print(f"\n[DEBUG] RAW RESPONSE:\n{response_text}\n")
        
        # 3. Parse Response
        slides = []
        current_slide = None
        
        lines = response_text.split('\n')
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # Detect Slide Title (e.g., "Slide 1: Title" or just "Title")
            if line.lower().startswith("slide") and ":" in line:
                if current_slide:
                    slides.append(current_slide)
                title_text = line.split(":", 1)[1].strip()
                current_slide = {"title": title_text, "content": []}
            elif line.startswith("-") or line.startswith("*"):
                if current_slide:
                    current_slide["content"].append(line.lstrip("-* ").strip())
            else:
                # If it looks like a title but doesn't have "Slide X:" prefix, we can treat it as a new slide if we have no content yet,
                # or just append to previous content. For safety with small models, let's treat bare lines as content unless they are very short and look like titles?
                # Actually, simpler: If we have no slide, start one.
                if not current_slide:
                    current_slide = {"title": line, "content": []}
                else:
                    # Treat as content
                    current_slide["content"].append(line)
        
        if current_slide:
            slides.append(current_slide)

        # Fallback if parsing failed completely
        if not slides:
             slides = [{"title": "Presentation", "content": [line for line in lines if line]}]

        # 4. Create PPTX
        for slide_info in slides:
            slide_layout = self.prs.slide_layouts[1] # Title and Content
            slide = self.prs.slides.add_slide(slide_layout)
            
            # Title
            title = slide.shapes.title
            title.text = slide_info.get("title", "Untitled")
            
            # Content
            content_shape = slide.placeholders[1]
            tf = content_shape.text_frame
            tf.clear()
            
            for point in slide_info.get("content", []):
                p = tf.add_paragraph()
                p.text = point
                p.font.size = Pt(18)
            
            # Add User Image if provided
            if image_path and os.path.exists(image_path):
                try:
                    self.prs.slides[self.prs.slides.index(slide)].shapes.add_picture(
                        image_path, Inches(7.5), Inches(5.5), width=Inches(2)
                    )
                except Exception as e:
                    print(f"Could not add image: {e}")

        # 5. Add Disclaimer Slide
        blank_layout = self.prs.slide_layouts[6] 
        end_slide = self.prs.slides.add_slide(blank_layout)
        
        txBox = end_slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = "Facut cu iubire de catre gPPTX, ATENTIE gPPTX POATE AVEA ERORI, SE RECOMANDA VERIFICAREA FISIERULUI FINAL!"
        p.font.bold = True
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER

        # 6. Save
        try:
            self.prs.save(output_file)
            return f"Presentation saved successfully to {output_file}"
        except Exception as e:
            return f"Error saving file: {e}"
